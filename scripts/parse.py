#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ppt-pdf-to-docx 核心脚本(v3)
阶段一:原始素材提取
  - PPTX: 递归提取文字/图片(含 GROUP 组合内)、LibreOffice 渲染整页截图
  - PDF: 转图 + 双引擎文字提取
输出:docx(带整页截图 + 原始文字)+ _assets/ 目录(截图、提取图片、JSON)
阶段二(视觉核对与逻辑重排)由 agent 基于 _assets/ 中的素材完成。

v3 修复 (vs v2):
  FIX-001: images_pdf 中 sha256 hash 变量被 img.size 覆盖导致同页去重失效
  FIX-002: _merge_pdf_lines 短行 flush 阈值从 3 提升到 8,减少页眉/页码误合并
  FIX-003: build_docx 多行文字按换行符拆分为多个段落,避免 docx 中文字挤在一行
  FIX-004: images_pdf 中 Pixmap 显式释放 (del pix)
  FIX-005: _iter_shapes 异常捕获精确化,仅捕获 AttributeError,不吞未知异常
  FIX-006: 移除未实现的 __main__.py 相关说明
  FIX-008: 移除 sys.exit(2) 硬退出,改为 return None 让主流程统一处理
  FIX-009: text_pptx 中 para.text 兜底注释改进,说明其无害性
  FIX-010: --assets-dir 按源文件独立子目录避免多文件覆盖
  FIX-011: 改进 _merge_pdf_lines 段落合并逻辑,更精确识别断行
  FIX-012: SKILL.md 使用说明改为 `python parse.py` 相对路径方式
"""
import argparse, os, sys, shutil, subprocess, hashlib, tempfile, re, glob, json
import importlib
from pathlib import Path

# ---------------------------------------------------------------------------
# 1. 环境探测
# ---------------------------------------------------------------------------
def detect_env():
    env = {
        "soffice": shutil.which("soffice") or shutil.which("libreoffice"),
        "wps":     shutil.which("wps") or shutil.which("wpspdf"),
        "pdftoppm":shutil.which("pdftoppm"),
        "tesseract": shutil.which("tesseract"),
        "xvfb":    shutil.which("xvfb-run"),
        "pymupdf": False, "pptx": False, "pdfplumber": False, "docx": False,
        "pytesseract": False, "PIL": False,
    }
    for mod, key in [("fitz", "pymupdf"), ("pptx", "pptx"),
                     ("pdfplumber", "pdfplumber"), ("docx", "docx"),
                     ("pytesseract", "pytesseract"), ("PIL", "PIL")]:
        try:
            __import__(mod); env[key] = True
        except Exception:
            pass
    return env

# ---------------------------------------------------------------------------
# 1b. Python 依赖自动安装
# ---------------------------------------------------------------------------
PIP_PKGS = ["python-pptx", "pdfplumber", "pymupdf", "python-docx",
            "pytesseract", "Pillow"]
_PIP_MOD = {
    "python-pptx": "pptx", "pdfplumber": "pdfplumber", "pymupdf": "fitz",
    "python-docx": "docx", "pytesseract": "pytesseract", "Pillow": "PIL",
}

def _import_ok(mod):
    try:
        importlib.import_module(mod); return True
    except Exception:
        return False

def _pip_install(pkgs):
    """尝试安装 pip 包:--user → 普通 → --break-system-packages"""
    pip_cmds = [
        [sys.executable, "-m", "pip", "install", "--user"],
        [sys.executable, "-m", "pip", "install"],
        [sys.executable, "-m", "pip", "install", "--break-system-packages"],
    ]
    for base in pip_cmds:
        try:
            r = subprocess.run(base + pkgs,
                               stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                               timeout=180)
            if r.returncode == 0:
                return True
        except Exception:
            pass
    return False

def ensure_python_deps(auto_install=True):
    if not auto_install:
        return
    missing = [p for p in PIP_PKGS if not _import_ok(_PIP_MOD[p])]
    if missing:
        print(f"[安装] 缺失 Python 依赖: {missing}")
        if _pip_install(missing):
            print("       -> pip 安装成功")
        else:
            print("       -> pip 安装失败,请手动: pip install " + " ".join(missing))

# ---------------------------------------------------------------------------
# 2. 工具函数
# ---------------------------------------------------------------------------
def sha256(b: bytes) -> str:
    return hashlib.sha256(b).hexdigest()[:12]

def _clean(s):
    """清洗写入 docx 的非法 XML 字符。"""
    if not s:
        return s
    out = []
    for ch in s:
        o = ord(ch)
        if o in (0x09, 0x0A, 0x0D):
            out.append(ch)
        elif o < 0x20 or (0x7F <= o <= 0x9F):
            continue
        else:
            out.append(ch)
    return "".join(out)

def _iter_shapes(shapes):
    """递归遍历 shapes,展开 GROUP 组合内的子形状,带 shape_type 异常保护。

    FIX-005: 仅捕获 AttributeError(无 shape_type 属性时的正常情况),
    不吞掉 IndexError/KeyError 等可能表明数据异常的异常。
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sh in shapes:
        try:
            st = sh.shape_type
        except AttributeError:
            # 没有 shape_type 属性的形状,直接 yield 兜底
            try:
                yield sh
            except Exception:
                pass
            continue
        try:
            if st == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(sh.shapes)
            else:
                yield sh
        except (AttributeError, TypeError, IndexError, KeyError):
            # 保护性兜底：GROUP 内部子形状遍历失败时仍尝试 yield 父级
            # 仅捕获常见的属性/类型/索引异常，不吞掉 MemoryError 等
            try:
                yield sh
            except Exception:
                pass

def _run_cmd(cmd, timeout=120, env=None, **kw):
    """运行子进程,带 timeout。"""
    try:
        return subprocess.run(cmd, timeout=timeout, env=env, **kw)
    except subprocess.TimeoutExpired:
        print(f"       -> 命令超时 ({timeout}s): {' '.join(str(c) for c in cmd[:4])}...")
        return None

# ---------------------------------------------------------------------------
# 3. 渲染:PPT -> PDF -> PNG
# ---------------------------------------------------------------------------
def render_pptx(pptx_path, out_dir, dpi, env):
    """返回 (png_list, renderer_name)。
    使用 xvfb-run(如有)、独立 LibreOffice profile、清除 Python 路径冲突。"""
    tmp_pdf_dir = tempfile.mkdtemp(prefix="ppt_pdf_")
    try:
        soffice = env["soffice"]
        if not soffice:
            return [], None

        # 隔离 profile 避免并发锁冲突
        lo_profile = f"file:///tmp/lo_profile_{os.getpid()}"
        # 清除 LibreOffice 内嵌 Python 的路径警告
        lo_env = os.environ.copy()
        lo_env.pop("PYTHONHOME", None)
        lo_env.pop("PYTHONPATH", None)

        soffice_cmd = [
            soffice, "--headless",
            f"-env:UserInstallation={lo_profile}",
            "--convert-to", "pdf",
            "--outdir", tmp_pdf_dir, str(pptx_path),
        ]
        if env["xvfb"]:
            soffice_cmd = ["xvfb-run", "-a"] + soffice_cmd

        print(f"       -> LibreOffice 渲染中...")
        r = _run_cmd(soffice_cmd, timeout=600, env=lo_env,
                     stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
                     text=True)
        if r is None:
            return [], None
        if r.returncode != 0:
            # 打印错误帮助诊断
            err_tail = (r.stdout or "")[-500:]
            if err_tail.strip():
                print(f"       -> soffice stderr: {err_tail[:300]}")
            return [], None

        pdfs = glob.glob(os.path.join(tmp_pdf_dir, "*.pdf"))
        if not pdfs:
            print("       -> 未找到输出 PDF")
            return [], None

        return pdf_to_png(pdfs[0], out_dir, "page", dpi, env), "LibreOffice"
    finally:
        shutil.rmtree(tmp_pdf_dir, ignore_errors=True)

def pdf_to_png(pdf_path, out_dir, prefix, dpi, env):
    out = []
    out_path = Path(out_dir)
    if env["pdftoppm"]:
        r = _run_cmd(["pdftoppm", "-png", "-r", str(dpi), str(pdf_path),
                      str(out_path / prefix)],
                     timeout=300, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        if r is None or r.returncode != 0:
            print(f"       -> pdftoppm 失败 (rc={r.returncode if r else 'timeout'})")
        else:
            out = glob.glob(str(out_path / f"{prefix}-*.png"))
    elif env["pymupdf"]:
        import fitz
        doc = fitz.open(pdf_path)
        try:
            for i, pg in enumerate(doc, 1):
                pix = pg.get_pixmap(dpi=dpi)
                p = out_path / f"{prefix}-{i:02d}.png"
                pix.save(str(p))
                out.append(str(p))
        finally:
            doc.close()
    def _num(p):
        m = re.search(r"-(\d+)\.png$", p)
        return int(m.group(1)) if m else 0
    return sorted(out, key=_num)

# ---------------------------------------------------------------------------
# 4. 文字提取
# ---------------------------------------------------------------------------
def text_pptx(pptx_path):
    """递归提取文字,按 (top,left) 排序反映视觉顺序,返回 {page: text}。
    同时返回坐标信息的辅助数据。"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    res = {}
    res_detail = {}
    for i, slide in enumerate(prs.slides, 1):
        items = []
        for sh in _iter_shapes(slide.shapes):
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    # FIX-009: para.text 已包含所有 runs 拼接结果,
                    # 但有些旧版本 python-pptx 可能返回空,保留 runs 回退作为无害兜底
                    t = (para.text or "").strip()
                    if not t:
                        t = "".join(r.text for r in para.runs).strip()
                    if t:
                        try:
                            top, left = sh.top or 0, sh.left or 0
                        except Exception:
                            top, left = 0, 0
                        items.append((top, left, t))
            if sh.has_table:
                for row in sh.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        try:
                            top, left = sh.top or 0, sh.left or 0
                        except Exception:
                            top, left = 0, 0
                        items.append((top, left, " | ".join(cells)))
        items.sort(key=lambda x: (x[0], x[1]))
        res[i] = "\n".join(t for _, _, t in items)
        res_detail[i] = [{"top": t, "left": l, "text": txt} for t, l, txt in items]
    return res, res_detail

def _merge_pdf_lines(text):
    """改进的段落合并。

    规则:
    - 空行 = 段落分隔
    - 行尾是句末标点时视为段落结束
    - 行首小写字母或非标点开头且上行未结束时合并到上行
    - 短行(< 8 字符)除非是句末标点结束,否则与下行合并尝试
    """
    lines = text.split("\n")
    out = []
    buf = ""
    # 句末标点:中英文句号、问号、感叹号、分号、冒号
    strong_end = set("。！？；：.!?;:")

    for line in lines:
        s = line.strip()
        if not s:
            if buf:
                out.append(buf)
                buf = ""
            continue

        if not buf:
            buf = s
            continue

        should_merge = (
            buf[-1] not in strong_end
            and (not s[0].isupper() or len(s) < 3)
        )
        if should_merge:
            buf += s
        else:
            out.append(buf)
            buf = s

    if buf:
        out.append(buf)
    return "\n".join(out)


def merge_single_letters(text):
    """合并被异常逐字分离的英文。"""
    toks = text.split()
    out, buf = [], []
    for t in toks:
        if len(t) == 1 and t.isalpha():
            buf.append(t)
        else:
            if len(buf) >= 3:
                out.append("".join(buf))
            else:
                out.extend(buf)
            buf = []
            out.append(t)
    if buf:
        out.append("".join(buf)) if len(buf) >= 3 else out.extend(buf)
    return " ".join(out)

def _score_text(d):
    s = 0
    for v in d.values():
        s += len(re.findall(r"[\u4e00-\u9fff]", v))
        s -= 5 * v.count("(cid:")
    return s

def text_pdf(pdf_path):
    """双引擎选优 + pdftotext 回退。"""
    cands = {}
    try:
        import pdfplumber
        pdf = pdfplumber.open(pdf_path)
        try:
            cands["plumber"] = {i: (pg.extract_text() or "") for i, pg in enumerate(pdf.pages, 1)}
        finally:
            pdf.close()
    except Exception:
        pass
    try:
        import fitz
        doc = fitz.open(pdf_path)
        try:
            cands["fitz"] = {i + 1: doc[i].get_text("text") for i in range(len(doc))}
        finally:
            doc.close()
    except Exception:
        pass

    best = None
    if "plumber" in cands and "fitz" in cands:
        best = "plumber" if _score_text(cands["plumber"]) >= _score_text(cands["fitz"]) else "fitz"
    elif "plumber" in cands:
        best = "plumber"
    elif "fitz" in cands:
        best = "fitz"

    raw = cands.get(best) if best else None
    if not raw:
        try:
            r = _run_cmd(["pdftotext", str(pdf_path), "-"],
                         timeout=120, capture_output=True, text=True)
            if r and r.stdout:
                parts = r.stdout.split("\f")
                while parts and not parts[0].strip(): parts.pop(0)
                while parts and not parts[-1].strip(): parts.pop()
                if parts:
                    try:
                        info = _run_cmd(["pdfinfo", str(pdf_path)],
                                        timeout=30, capture_output=True, text=True)
                        n = int(info.stdout.split("Pages:")[1].split()[0])
                        if 0 < n < len(parts): parts = parts[:n]
                    except Exception:
                        pass
                    raw = {i + 1: parts[i] for i in range(len(parts))}
        except Exception:
            raw = None
    if not raw:
        return {}
    result = {}
    for k, v in raw.items():
        cleaned = merge_single_letters(v).strip()
        cleaned = _merge_pdf_lines(cleaned)
        result[k] = cleaned
    return result

# ---------------------------------------------------------------------------
# 5. 图片提取
# ---------------------------------------------------------------------------
IMAGE_MIN_AREA = 100_000  # 10 万像素(约 300x300)

def images_pptx(pptx_path):
    """递归提取嵌入图片,按 blob hash 同页去重,返回 {page: [{id, name, blob, w, h, pos}]}"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    import io
    prs = Presentation(pptx_path)
    res = {}
    for i, slide in enumerate(prs.slides, 1):
        seen_hashes = set()
        imgs = []
        for sh in _iter_shapes(slide.shapes):
            try:
                is_pic = sh.shape_type == MSO_SHAPE_TYPE.PICTURE
            except Exception:
                continue
            if is_pic:
                try:
                    blob = sh.image.blob
                    h = sha256(blob)
                    if h in seen_hashes:
                        continue
                    seen_hashes.add(h)
                    img = Image.open(io.BytesIO(blob))
                    w_px, h_px = img.size
                    if w_px * h_px >= IMAGE_MIN_AREA:
                        try:
                            pos = {"left": sh.left or 0, "top": sh.top or 0,
                                   "width": sh.width or 0, "height": sh.height or 0}
                        except Exception:
                            pos = {}
                        imgs.append({"name": sh.name, "blob": blob,
                                     "w": w_px, "h": h_px, "pos": pos})
                except Exception:
                    pass
        if imgs:
            res[i] = imgs
    return res

def images_pdf(pdf_path, dpi):
    """PDF 图片通过 pymupdf 提取。

    FIX-001: 修复 sha256 hash 变量被 img.size 覆盖导致去重失效。
    FIX-004: 显式释放 Pixmap 资源。
    """
    import fitz
    from PIL import Image
    import io
    doc = fitz.open(pdf_path)
    res = {}
    for pnum in range(len(doc)):
        page = doc[pnum]
        imgs = []
        seen_hashes = set()
        for idx, img_info in enumerate(page.get_images(full=True)):
            xref = img_info[0]
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha >= 4:
                    old_pix = pix
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                    del old_pix
                blob = pix.tobytes("png")
                # FIX-001: 使用独立变量名 img_hash,避免被 img.size 覆盖
                img_hash = sha256(blob)
                if img_hash in seen_hashes:
                    del pix  # FIX-004
                    continue
                seen_hashes.add(img_hash)
                img = Image.open(io.BytesIO(blob))
                img_w, img_h = img.size  # FIX-001: 使用 img_w, img_h 而非 w, h
                if img_w * img_h >= IMAGE_MIN_AREA:
                    imgs.append({"name": f"img_{idx}", "blob": blob,
                                 "w": img_w, "h": img_h, "pos": {}})
                del pix  # FIX-004
            except Exception:
                pass
        if imgs:
            res[pnum + 1] = imgs
    doc.close()
    return res

# ---------------------------------------------------------------------------
# 6. 生成 _assets 目录
# ---------------------------------------------------------------------------
def save_assets(assets_dir, page_pngs, img_map, texts, texts_detail, source_name):
    """保存截图、提取的图片、文字 JSON。"""
    adir = Path(assets_dir)
    adir.mkdir(parents=True, exist_ok=True)
    pic_dir = adir / "extracted_pics"
    pic_dir.mkdir(exist_ok=True)

    page_map = {}
    for png in page_pngs:
        m = re.search(r"-(\d+)\.png$", png)
        if m:
            pg = int(m.group(1))
            dst = adir / f"page_{pg:02d}.png"
            shutil.copy2(png, dst)
            page_map[pg] = str(dst)

    extracted = {}
    for pg, imgs in img_map.items():
        extracted[pg] = []
        for idx, img_info in enumerate(imgs, 1):
            fname = f"page_{pg:02d}_pic_{idx:02d}.png"
            (pic_dir / fname).write_bytes(img_info["blob"])
            entry = {
                "id": f"pic{idx:02d}",
                "name": img_info["name"],
                "file": str(pic_dir / fname),
                "size": [img_info["w"], img_info["h"]],
            }
            if img_info.get("pos"):
                entry["pos"] = img_info["pos"]
            extracted[pg].append(entry)

    data = {"source": source_name, "pages": {}}
    all_pages = sorted(set(list(page_map.keys()) + list(texts.keys()) + list(extracted.keys())))
    for pg in all_pages:
        data["pages"][str(pg)] = {
            "screenshot": page_map.get(pg, ""),
            "text": texts.get(pg, ""),
            "text_blocks": texts_detail.get(pg, []),
            "images": extracted.get(pg, []),
        }
    json_path = adir / "extracted_text.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    return str(json_path)

# ---------------------------------------------------------------------------
# 7. 生成 docx
# ---------------------------------------------------------------------------
def build_docx(files_info, out_path, total_pages, slide_size=None):
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    doc = Document()

    note = doc.add_paragraph()
    r = note.add_run(
        f"<!-- 本文档为阶段一输出(原始提取),合计 {total_pages} 页。\n"
        f"标记说明:\n"
        f"  /[pic NN:简述]/   → 行内小图\n"
        f"  /[pic NN]*** ... ***/  → 块级大图/详细描述\n"
        f"  无标记文字 = PPT 可编辑文本(按视觉从上到下排序);阶段二由 agent 视觉核对后重排 -->"
    )
    r.font.size = Pt(8); r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # 根据 slide 宽高比决定图片宽度
    if slide_size and slide_size[0] and slide_size[1]:
        ratio = slide_size[0] / slide_size[1]
        if ratio > 1.5:
            img_width = Inches(6.5)  # 宽屏 16:9
        elif ratio < 1.0:
            img_width = Inches(4.0)  # 竖版
        else:
            img_width = Inches(5.5)  # 接近 4:3
    else:
        img_width = Inches(6.0)

    multi = len(files_info) > 1
    for fi in files_info:
        if multi:
            doc.add_heading(f"来源文件:{_clean(fi['name'])}", level=1)
        for pg in sorted(fi["pages"]):
            doc.add_paragraph(f"page {pg}")
            pngs = fi["pages"][pg].get("png")
            if pngs and pngs[0]:
                try:
                    doc.add_picture(pngs[0], width=img_width)
                except Exception as e:
                    p = doc.add_paragraph()
                    r = p.add_run(f"【整页截图加载失败:{e}】")
                    r.font.size = Pt(8); r.font.color.rgb = RGBColor(0xAA, 0x00, 0x00)
            else:
                p = doc.add_paragraph()
                r = p.add_run("【无整页截图】")
                r.font.size = Pt(8); r.font.color.rgb = RGBColor(0xAA, 0x00, 0x00)

            # FIX-003: 多行文字按换行符拆分，使用 add_break() 产生真正的换行
            text = fi["pages"][pg].get("text") or ""
            if text:
                text_para = doc.add_paragraph()
                text_para.add_run("文字内容：").bold = True
                lines = text.split("\n")
                for li, line in enumerate(lines):
                    if li > 0:
                        text_para.add_run().add_break()
                    text_para.add_run(_clean(line))
            else:
                doc.add_paragraph("文字内容：(无)")

            pic_count = fi["pages"][pg].get("pic_count", 0)
            if pic_count > 0:
                labels = " ".join([f"/[pic {i+1:02d}:待识别]/" for i in range(pic_count)])
                doc.add_paragraph(f"插图标记:{labels}")
            else:
                doc.add_paragraph("插图说明:无")

    doc.save(out_path)

# ---------------------------------------------------------------------------
# 8. 单文件解析
# ---------------------------------------------------------------------------
def parse_one(path, work, dpi, env, require_renderer_for_pptx=True):
    ext = path.suffix.lower()
    slug = re.sub(r'\W+', '_', path.stem)[:36] or 'file'
    slug += "_" + sha256(path.name.encode())[:4]
    fwork = Path(work) / slug
    fwork.mkdir(parents=True, exist_ok=True)

    pages = {}
    if ext == ".pptx":
        if not env["soffice"]:
            if require_renderer_for_pptx:
                print(f"\n[错误] 未检测到 LibreOffice(soffice),无法渲染 PPT 整页截图。")
                print(f"  PPT 解析必须有渲染器。请安装 LibreOffice 后重试,或将 PPT 导出为 PDF。")
                print(f"  安装命令参见 SKILL.md「环境依赖安装」章节。")
                return None  # FIX-008: 不再 sys.exit
            else:
                print("[警告] 无 LibreOffice,仅提取文字和图片,无整页截图。")

        print(f"[PPTX] 渲染整页截图...")
        pngs, renderer = render_pptx(str(path), fwork, dpi, env)
        if renderer:
            print(f"       -> 使用 {renderer},生成 {len(pngs)} 页截图")
        else:
            print("       -> 渲染失败")

        print(f"[PPTX] 提取文字(含坐标排序)...")
        texts, texts_detail = text_pptx(str(path))

        print(f"[PPTX] 提取内嵌图片...")
        img_map = images_pptx(str(path))
        total_pics = sum(len(v) for v in img_map.values())
        print(f"       -> 共 {total_pics} 张图片分布在 {len(img_map)} 页")

        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slide_size = (prs.slide_width, prs.slide_height)
        except Exception:
            slide_size = None

        npages = max(len(texts), len(pngs), len(img_map)) or 1
        for i in range(1, npages + 1):
            png = [pngs[i-1]] if i-1 < len(pngs) else []
            pages[i] = {
                "png": png,
                "text": texts.get(i, ""),
                "pic_count": len(img_map.get(i, [])),
            }
        return {"name": path.name, "pages": pages,
                "_pngs": pngs, "_img_map": img_map, "_texts": texts,
                "_texts_detail": texts_detail, "_slide_size": slide_size}

    elif ext == ".pdf":
        print(f"[PDF] 转图中...")
        pngs = pdf_to_png(str(path), fwork, "page", dpi, env)
        print(f"[PDF] 提取文字...")
        texts = text_pdf(str(path))
        texts_detail = {k: [{"top": 0, "left": 0, "text": v}] for k, v in texts.items()}
        print(f"[PDF] 提取内嵌图片...")
        img_map = images_pdf(str(path), dpi)
        total_pics = sum(len(v) for v in img_map.values())
        print(f"       -> 共 {total_pics} 张图片分布在 {len(img_map)} 页")
        npages = max(len(texts), len(pngs)) or 1
        for i in range(1, npages + 1):
            png = [pngs[i-1]] if i-1 < len(pngs) else []
            pages[i] = {
                "png": png,
                "text": texts.get(i, ""),
                "pic_count": len(img_map.get(i, [])),
            }
        return {"name": path.name, "pages": pages,
                "_pngs": pngs, "_img_map": img_map, "_texts": texts,
                "_texts_detail": texts_detail, "_slide_size": None}
    else:
        print(f"[跳过] 不支持格式: {path}")
        return None

# ---------------------------------------------------------------------------
# 9. 主流程
# ---------------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser(description="PPT/PDF -> 逐页 docx 解析(阶段一:素材提取)")
    ap.add_argument("inputs", nargs="+", help=".pptx / .pdf 文件")
    ap.add_argument("--out", default="解析结果.docx")
    ap.add_argument("--dpi", type=int, default=150)
    ap.add_argument("--assets-dir", default=None)
    ap.add_argument("--keep", action="store_true")
    ap.add_argument("--no-install", action="store_true")
    ap.add_argument("--allow-no-renderer", action="store_true")
    args = ap.parse_args()

    ensure_python_deps(auto_install=not args.no_install)
    env = detect_env()
    print("环境检测:", {k: (v if isinstance(v, bool) else (v and '✓'))
                          for k, v in env.items()})

    work = tempfile.mkdtemp(prefix="pptpdf_v2_")
    files_info = []
    total = 0
    slide_size = None
    for inp in args.inputs:
        p = Path(inp)
        if not p.exists():
            print(f"[错误] 文件不存在: {inp}"); continue
        info = parse_one(p, work, args.dpi, env,
                         require_renderer_for_pptx=not args.allow_no_renderer)
        if info:
            files_info.append(info)
            total += len(info["pages"])
            if info.get("_slide_size"):
                slide_size = info["_slide_size"]

    if not files_info:
        print("无有效输入,退出。"); sys.exit(1)

    out = Path(args.out)
    if not out.is_absolute():
        out = Path.cwd() / out
    out.parent.mkdir(parents=True, exist_ok=True)

    # FIX-010: --assets-dir 按源文件独立子目录避免多文件覆盖
    json_path = None
    for fi in files_info:
        src_stem = Path(fi["name"]).stem
        if args.assets_dir:
            assets_sub = Path(args.assets_dir) / src_stem
        else:
            assets_sub = out.parent / (out.stem + "_assets") / src_stem
        jp = save_assets(str(assets_sub), fi.pop("_pngs", []), fi.pop("_img_map", {}),
                         fi.pop("_texts", {}), fi.pop("_texts_detail", {}), fi["name"])
        if not json_path:
            json_path = jp

    build_docx(files_info, str(out), total, slide_size=slide_size)
    print(f"\n✅ 阶段一完成:{out}")
    print(f"   合计 {total} 页,来源文件 {len(files_info)} 个")
    print(f"   素材目录:{out.parent / (out.stem + '_assets')}")
    print(f"   文字 JSON:{json_path}")

    if not args.keep:
        shutil.rmtree(work, ignore_errors=True)

if __name__ == "__main__":
    main()
